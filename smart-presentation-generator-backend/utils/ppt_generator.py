# utils/ppt_generator.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os
from models.templates import SLIDE_TEMPLATES

class PPTGenerator:
    """Utility class for generating PowerPoint presentations."""
    
    def __init__(self, template_path=None):
        """Initialize with optional template."""
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
    
    def add_title_slide(self, title, subtitle=None, theme_settings=None):
        """Add a title slide to the presentation."""
        slide_layout = self.prs.slide_layouts[SLIDE_TEMPLATES["title_slide"].layout_type]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Add subtitle if provided
        if subtitle and hasattr(slide, 'placeholders') and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
        
        # Apply theme if provided
        if theme_settings:
            self._apply_theme_to_slide(slide, theme_settings, is_title_slide=True)
        
        return slide
    
    def add_content_slide(self, title, content, layout_type=1, theme_settings=None):
        """Add a content slide with bullet points."""
        slide_layout = self.prs.slide_layouts[layout_type]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Add content as bullet points
        if hasattr(slide, 'placeholders') and len(slide.placeholders) > 1:
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            
            # Clear any existing text
            for paragraph in text_frame.paragraphs:
                if paragraph.text:
                    paragraph.text = ""
            
            # Add new content
            first_paragraph = True
            for point in content:
                if first_paragraph and text_frame.paragraphs:
                    p = text_frame.paragraphs[0]
                    first_paragraph = False
                else:
                    p = text_frame.add_paragraph()
                p.text = point
                p.level = 0  # First level bullet
        
        # Apply theme if provided
        if theme_settings:
            self._apply_theme_to_slide(slide, theme_settings)
        
        return slide
    
    def add_two_column_slide(self, title, left_content, right_content, theme_settings=None):
        """Add a slide with two columns of content."""
        slide_layout = self.prs.slide_layouts[SLIDE_TEMPLATES["two_content"].layout_type]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Add left column content
        if hasattr(slide, 'placeholders') and len(slide.placeholders) > 1:
            left_shape = slide.placeholders[1]
            text_frame = left_shape.text_frame
            
            for point in left_content:
                p = text_frame.add_paragraph()
                p.text = point
                p.level = 0
        
        # Add right column content
        if hasattr(slide, 'placeholders') and len(slide.placeholders) > 2:
            right_shape = slide.placeholders[2]
            text_frame = right_shape.text_frame
            
            for point in right_content:
                p = text_frame.add_paragraph()
                p.text = point
                p.level = 0
        
        # Apply theme if provided
        if theme_settings:
            self._apply_theme_to_slide(slide, theme_settings)
        
        return slide
    
    def _apply_theme_to_slide(self, slide, theme_settings, is_title_slide=False):
        """Apply theme settings to a slide."""
        # Background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*theme_settings['background_color'])
        
        # Text formatting for title
        title_shape = slide.shapes.title
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(theme_settings['title_font_size'])
            paragraph.font.name = theme_settings['title_font']
            paragraph.font.color.rgb = RGBColor(*theme_settings['title_color'])
        
        # Additional formatting for content if not a title slide
        if not is_title_slide:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape != title_shape:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.size = Pt(theme_settings['content_font_size'])
                        paragraph.font.name = theme_settings['content_font']
                        paragraph.font.color.rgb = RGBColor(*theme_settings['content_color'])
    
    def save(self, file_path):
        """Save the presentation to the specified path."""
        # Create directory if it doesn't exist
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        self.prs.save(file_path)
        return file_path